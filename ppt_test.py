from pptx import Presentation
import os 

path=os.path.dirname(__file__)

prs = Presentation()

#使用ppt自带模板，ppt自带了常用的1-48种模板通过index选择对应的模板
title_slide_layout = prs.slide_layouts[0]

#新建一页幻灯片
slide = prs.slides.add_slide(title_slide_layout)

#获取幻灯片中的title元素(本页幻灯片必须含有标题元素才能通过此方法获取)
title = slide.shapes.title

#根据placeholdes索引获取一页幻灯片中的元素
subtitle = slide.shapes.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "pip install python-pptx"

prs.save(path+"/test.pptx")
