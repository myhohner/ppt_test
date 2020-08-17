from pptx import Presentation
from pptx.util import Inches
import os 
path=os.path.dirname(__file__)

prs=Presentation()
slide=prs.slides.add_slide(prs.slide_layouts[0])
#插入picture
'''
left,top,width,height = Inches(1),Inches(1),Inches(5),Inches(5)
pic_path=os.path.join(path,'2.jpg')
slide.shapes.add_picture(pic_path,left,top,width,height)
'''

#插入文本框
left, top, width, height = Inches(5),Inches(4),Inches(8),Inches(8)
tf = slide.shapes.add_textbox(left=left,top=top,width=width,height=height).text_frame
tf.text = "第一段"  #插入第一段文本
tf.text = "第二段"  #插入第二段文本
tf.text = "第三段"  #插入第三段文本
	#此时只会显示出"第三段"内容

left, top, width, height = Inches(3),Inches(6),Inches(8),Inches(2)
	#重新设定参数后，会有新的文本框生成
tf = slide.shapes.add_textbox(left=left,top=top,width=width,height=height).text_frame
tf.add_paragraph()#插入空白段
tf.paragraphs[0].text = "第二段"  #插入第二段文本
tf.paragraphs[1].text = "第三段"  #插入第二段文本
prs.save(path+'/pic&test.pptx')