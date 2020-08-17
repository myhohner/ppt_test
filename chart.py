from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import os 
path=os.path.dirname(__file__)
# create presentation with 1 slide ------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# define chart data ---------------------
chart_data = ChartData()
chart_data.categories = ['East', 'West', 'Midwest']
chart_data.add_series('Series 1', (19.2, 21.4, 16.7))

# add chart to slide --------------------COLUMN_CLUSTERED为柱状图，x,y为图形左上角点所在坐标，cx,cy为图形长、宽
x, y, cx, cy = Inches(2), Inches(2), Inches(5), Inches(4.5)
slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)


prs.save(path+'/chart-01.pptx')

